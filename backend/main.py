from fileinput import filename

from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from dotenv import load_dotenv
from groq import Groq

import os
import uuid
import requests
import base64
import urllib.parse
import PyPDF2
from pptx import Presentation

load_dotenv()

# =========================
# SETUP
# =========================
app = FastAPI()

os.makedirs("generated_images", exist_ok=True)
app.mount(
    "/generated_images",
    StaticFiles(directory="generated_images"),
    name="generated_images"
)

# =========================
# CORS
# =========================
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# =========================
# ENV / KEYS
# =========================
groq_key = os.getenv("GROQ_API_KEY")
cloudflare_token = os.getenv("CLOUDFLARE_API_TOKEN")
cloudflare_account_id = os.getenv("CLOUDFLARE_ACCOUNT_ID")

print("GROQ KEY FOUND:", bool(groq_key))
print("CLOUDFLARE TOKEN FOUND:", bool(cloudflare_token))
print("CLOUDFLARE ACCOUNT ID FOUND:", bool(cloudflare_account_id))

client = Groq(api_key=groq_key)


# =========================
# REQUEST MODEL
# =========================
class ChatRequest(BaseModel):
    message: str
    fileContext: str | None = ""


# =========================
# HELPERS
# =========================
def clean_image_prompt(prompt: str) -> str:
    cleaned = prompt.lower().strip()

    remove_phrases = [
        "generate image of",
        "generate an image of",
        "generate image",
        "create image of",
        "create an image of",
        "create image",
        "make image of",
        "make an image of",
        "make image",
        "draw",
        "illustrate",
        "show me",
        "image of",
        "picture of",
        "photo of"
    ]

    for phrase in remove_phrases:
        if cleaned.startswith(phrase):
            cleaned = cleaned.replace(phrase, "", 1).strip()
            break

    return cleaned or "beautiful futuristic digital art"


def pollinations_url(prompt: str) -> str:
    encoded = urllib.parse.quote(prompt)
    return f"https://image.pollinations.ai/prompt/{encoded}"


# =========================
# FILE EXTRACTION API
# =========================
@app.post("/extract")
async def extract_file(file: UploadFile = File(...)):
    try:
        text = ""
        filename = (file.filename or "").lower()

        if filename.endswith(".pdf"):
            pdf = PyPDF2.PdfReader(file.file)
            for page in pdf.pages:
                text += (page.extract_text() or "") + "\n"

        elif filename.endswith(".pptx"):
            prs = Presentation(file.file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"

        else:
            return {"text": ""}

        print("EXTRACTED LENGTH:", len(text))
        return {"text": text.strip()}

    except Exception as e:
        print("EXTRACT ERROR:", str(e))
        return {"text": ""}


# =========================
# IMAGE GENERATION API
# =========================
@app.post("/generate-image")
def generate_image(req: ChatRequest):
    try:
        prompt = req.message.strip()

        if not prompt:
            return {"reply": "Prompt is empty", "type": "text"}

        cleaned = clean_image_prompt(prompt)

        # fallback if cloudflare not configured
        if not cloudflare_token or not cloudflare_account_id:
            return {
                "reply": pollinations_url(cleaned),
                "type": "image"
            }

        unique_id = uuid.uuid4().hex[:12]

        url = f"https://api.cloudflare.com/client/v4/accounts/{cloudflare_account_id}/ai/run/@cf/black-forest-labs/flux-1-schnell"

        headers = {
            "Authorization": f"Bearer {cloudflare_token}",
            "Content-Type": "application/json"
        }

        payload = {
            "prompt": cleaned,
            "steps": 6
        }

        print("IMAGE PROMPT:", cleaned)
        print("CALLING CLOUDFLARE...")

        response = requests.post(url, headers=headers, json=payload, timeout=90)

        if response.status_code != 200:
            print("CLOUDFLARE ERROR:", response.text)
            return {
                "reply": pollinations_url(cleaned),
                "type": "image"
            }

        data = response.json()
        image_b64 = data.get("result", {}).get("image")

        if not image_b64:
            print("INVALID CLOUDFLARE RESPONSE:", data)
            return {
                "reply": pollinations_url(cleaned),
                "type": "image"
            }

        image_bytes = base64.b64decode(image_b64)

        filename = f"{unique_id}.jpg"
        filepath = os.path.join("generated_images", filename)

        with open(filepath, "wb") as f:
            f.write(image_bytes)

        local_url = f"https://vibebot-backend.onrender.com/generated_images/{filename}"

        return {
            "reply": local_url,
            "type": "image"
        }

    except Exception as e:
        print("IMAGE ERROR:", str(e))
        return {
            "reply": pollinations_url(clean_image_prompt(req.message)),
            "type": "image"
        }


# =========================
# CHAT API
# =========================
@app.post("/chat")
def chat(req: ChatRequest):
    try:
        msg = req.message.strip()
        file_context = (req.fileContext or "").strip()

        print("\n------ NEW REQUEST ------")
        print("MESSAGE:", msg)
        print("FILE LENGTH:", len(file_context))

        quick_responses = {
            "hi": "Hello 👋",
            "hello": "Hello 👋 How can I help you?",
            "hey": "Hey 👋",
            "yo": "Yo 👋"
        }

        if msg.lower() in quick_responses:
            return {
                "reply": quick_responses[msg.lower()],
                "type": "text"
            }

        has_document = len(file_context.split()) > 20

        if has_document:
            prompt = f"""
You are a document assistant.

Answer ONLY using the document below.

DOCUMENT:
{file_context[:4000]}

QUESTION:
{msg}

If the answer is not present in the document, reply exactly:
Not found in document.
"""
        else:
            prompt = msg

        response = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[
                {"role": "user", "content": prompt}
            ],
            temperature=0.2,
            max_tokens=400
        )

        reply = response.choices[0].message.content or "No response generated."

        return {
            "reply": reply,
            "type": "text"
        }

    except Exception as e:
        print("CHAT ERROR:", str(e))
        return {
            "reply": f"ERROR: {str(e)}",
            "type": "text"
        }